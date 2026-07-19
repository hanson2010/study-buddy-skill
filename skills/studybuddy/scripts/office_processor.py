import os
import sys
import shutil

try:
    from docx import Document
except ImportError:
    print('Error: python-docx is not installed. Please install it with "pip install python-docx".', file=sys.stderr)
    sys.exit(1)

try:
    import openpyxl
except ImportError:
    print('Error: openpyxl is not installed. Please install it with "pip install openpyxl".', file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print('Error: python-pptx is not installed. Please install it with "pip install python-pptx".', file=sys.stderr)
    sys.exit(1)

try:
    import xlrd
except ImportError:
    print('Error: xlrd is not installed. Please install it with "pip install xlrd".', file=sys.stderr)
    sys.exit(1)

from common import get_data_dir, get_raw_dir, get_unique_filename, update_markdown_references

MAX_SIZE_MB = 50
MAX_SIZE_BYTES = MAX_SIZE_MB * 1024 * 1024
SUPPORTED_EXTENSIONS = ('.docx', '.xlsx', '.xls', '.pptx')
DEPRECATED_EXTENSIONS = ('.doc', '.ppt')


def extract_docx_text(filepath):
    try:
        doc = Document(filepath)
        paragraphs = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            style_name = para.style.name
            if style_name.startswith('Heading 1'):
                paragraphs.append('# ' + para.text)
            elif style_name.startswith('Heading 2'):
                paragraphs.append('## ' + para.text)
            elif style_name.startswith('Heading 3'):
                paragraphs.append('### ' + para.text)
            elif style_name.startswith('Heading 4'):
                paragraphs.append('#### ' + para.text)
            elif style_name.startswith('Heading 5'):
                paragraphs.append('##### ' + para.text)
            elif style_name.startswith('Heading 6'):
                paragraphs.append('###### ' + para.text)
            elif 'List' in style_name or 'Bullet' in style_name:
                paragraphs.append('- ' + para.text)
            elif 'Numbered' in style_name or 'Numbering' in style_name:
                paragraphs.append('1. ' + para.text)
            else:
                paragraphs.append(para.text)
        tables = []
        for table in doc.tables:
            md_table = []
            col_count = 0
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                col_count = len(cells)
                md_table.append('| ' + ' | '.join(cells) + ' |')
            if md_table and col_count > 0:
                md_table.insert(1, '| ' + ' | '.join(['---'] * col_count) + ' |')
                tables.append('\n'.join(md_table))
        content = '\n\n'.join(paragraphs)
        if tables:
            content += '\n\n' + '\n\n'.join(tables)
        return content.strip(), len(paragraphs)
    except Exception as e:
        print(f'Error extracting text from DOCX: {e}', file=sys.stderr)
        return None, 0


def extract_xlsx_text(filepath):
    try:
        wb = openpyxl.load_workbook(filepath, read_only=True)
        sheet_count = len(wb.sheetnames)
        sheets_content = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = []
                for cell in row:
                    if cell is None:
                        cells.append('')
                    else:
                        cells.append(str(cell))
                if any(c.strip() for c in cells):
                    rows.append(cells)
            if rows:
                md_table = []
                for i, row_data in enumerate(rows):
                    md_table.append('| ' + ' | '.join(row_data) + ' |')
                    if i == 0:
                        md_table.append('| ' + ' | '.join(['---'] * len(row_data)) + ' |')
                sheets_content.append(f'## {sheet_name}\n\n' + '\n'.join(md_table))
        wb.close()
        return '\n\n'.join(sheets_content).strip(), sheet_count
    except Exception as e:
        print(f'Error extracting text from XLSX: {e}', file=sys.stderr)
        return None, 0


def extract_xls_text(filepath):
    try:
        wb = xlrd.open_workbook(filepath)
        sheets_content = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            rows = []
            for row_idx in range(ws.nrows):
                cells = []
                for col_idx in range(ws.ncols):
                    cell = ws.cell_value(row_idx, col_idx)
                    if cell == '' or cell is None:
                        cells.append('')
                    else:
                        cells.append(str(cell))
                if any(c.strip() for c in cells):
                    rows.append(cells)
            if rows:
                md_table = []
                for i, row_data in enumerate(rows):
                    md_table.append('| ' + ' | '.join(row_data) + ' |')
                    if i == 0:
                        md_table.append('| ' + ' | '.join(['---'] * len(row_data)) + ' |')
                sheets_content.append(f'## {sheet_name}\n\n' + '\n'.join(md_table))
        return '\n\n'.join(sheets_content).strip(), len(wb.sheet_names())
    except Exception as e:
        print(f'Error extracting text from XLS: {e}', file=sys.stderr)
        return None, 0


def extract_pptx_text(filepath):
    try:
        prs = Presentation(filepath)
        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_runs.append(shape.text)
            if text_runs:
                slides_content.append(f'## Slide {i}\n\n' + '\n\n'.join(text_runs))
        return '\n\n'.join(slides_content).strip(), len(prs.slides)
    except Exception as e:
        print(f'Error extracting text from PPTX: {e}', file=sys.stderr)
        return None, 0


def process_docx(filepath):
    text, para_count = extract_docx_text(filepath)
    if text:
        return {'text': text, 'metadata': {'paragraphs': para_count}}
    return None


def process_xlsx(filepath):
    text, sheet_count = extract_xlsx_text(filepath)
    if text:
        return {'text': text, 'metadata': {'sheets': sheet_count}}
    return None


def process_xls(filepath):
    text, sheet_count = extract_xls_text(filepath)
    if text:
        return {'text': text, 'metadata': {'sheets': sheet_count}}
    return None


def process_pptx(filepath):
    text, slide_count = extract_pptx_text(filepath)
    if text:
        return {'text': text, 'metadata': {'slides': slide_count}}
    return None


def process_office(filepath):
    if not os.path.exists(filepath):
        print(f'Error: File not found - {filepath}', file=sys.stderr)
        return None
    
    file_ext = os.path.splitext(filepath)[1].lower()
    
    if file_ext in DEPRECATED_EXTENSIONS:
        print(f'Warning: {file_ext} format is not supported. Please convert to .docx or .pptx format first.')
        return None
    
    if file_ext not in SUPPORTED_EXTENSIONS:
        print(f'Skipping unsupported format - {filepath}')
        return None
    
    filesize = os.path.getsize(filepath)
    if filesize > MAX_SIZE_BYTES:
        print(f'Warning: File size ({filesize/1024/1024:.2f}MB) exceeds recommended limit of {MAX_SIZE_MB}MB')
    
    original_filename = os.path.basename(filepath)
    
    raw_dir = get_raw_dir()
    new_filename = get_unique_filename(raw_dir, original_filename)
    target_path = os.path.join(raw_dir, new_filename)
    
    print(f'Copying {filepath} ({filesize/1024/1024:.2f}MB) -> {target_path}')
    shutil.copy2(filepath, target_path)
    
    processor_funcs = {
        '.docx': process_docx,
        '.xlsx': process_xlsx,
        '.xls': process_xls,
        '.pptx': process_pptx,
    }
    
    result = processor_funcs.get(file_ext)(filepath)
    if result:
        text_filename = os.path.splitext(new_filename)[0] + '_text.md'
        text_path = os.path.join(raw_dir, text_filename)
        with open(text_path, 'w', encoding='utf-8') as f:
            f.write('---\n')
            f.write(f'office_file: {new_filename}\n')
            f.write(f'file_type: {file_ext[1:]}\n')
            for k, v in result['metadata'].items():
                f.write(f'{k}: {v}\n')
            f.write('---\n\n')
            f.write(result['text'])
        print(f'Extracted text saved to: {text_filename}')
    
    rel_path = os.path.relpath(target_path, get_data_dir())
    print(f'Result: {rel_path}')
    return rel_path


def main():
    if len(sys.argv) < 2:
        print('Usage: python office_processor.py <office_file> [<office_file> ...]')
        print('Supported formats: docx, xlsx, xls, pptx')
        print('Note: .doc and .ppt formats are not supported.')
        sys.exit(1)
    
    data_dir = get_data_dir()
    os.makedirs(data_dir, exist_ok=True)
    
    for filepath in sys.argv[1:]:
        result = process_office(filepath)
        if result:
            update_markdown_references(data_dir, filepath, result)


if __name__ == '__main__':
    main()
